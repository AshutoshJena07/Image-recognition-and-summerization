"""Universal File Parser for Images, Documents, Videos, Spreadsheets, and Text files."""

from __future__ import annotations

import io
import os
import re
from pathlib import Path
from typing import Any, Tuple

from PIL import Image

try:
    import pymupdf
except ImportError:
    try:
        import fitz as pymupdf
    except ImportError:
        pymupdf = None

try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import cv2
except ImportError:
    cv2 = None

# Known plain-text / code file extensions
TEXT_EXTENSIONS = {
    ".txt", ".json", ".md", ".py", ".js", ".ts", ".jsx", ".tsx",
    ".html", ".css", ".log", ".csv", ".xml", ".yaml", ".yml",
    ".c", ".cpp", ".h", ".sql", ".sh", ".bat", ".ps1"
}


def parse_uploaded_file(file_bytes: bytes, filename: str, content_type: str = "") -> Tuple[Image.Image | None, str]:
    """
    Process any uploaded file and return a tuple of (pil_image, extracted_text_context).
    - For images: returns (pil_image, "")
    - For documents (PDF, Word, PPT, Excel, Text, Code): returns (None, extracted_text)
    - For videos (MP4, AVI, MOV): extracts key frame as pil_image and returns (pil_image, video_metadata_text)
    """
    ext = Path(filename).suffix.lower()
    content_type = (content_type or "").lower()

    # 1. Image Formats (PNG, JPG, JPEG, WEBP, GIF, BMP, SVG)
    if ext in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif"} or "image/" in content_type:
        try:
            img = Image.open(io.BytesIO(file_bytes)).convert("RGB")
            return img, ""
        except Exception:
            pass

    if ext == ".svg":
        # Extract text/labels from SVG file
        svg_text = file_bytes.decode("utf-8", errors="ignore")
        clean_text = re.sub(r"<[^>]+>", " ", svg_text)
        clean_text = re.sub(r"\s+", " ", clean_text).strip()
        return None, f"[SVG Vector File Content & Text Elements]:\n{clean_text[:2000]}"

    # 2. PDF Documents (.pdf)
    if ext == ".pdf" or "pdf" in content_type:
        full_pdf_text = ""
        # Priority A: PyMuPDF (pymupdf) - High accuracy layout + image rendering
        if pymupdf:
            try:
                doc = pymupdf.open(stream=file_bytes, filetype="pdf")
                pages_text = []
                for i, page in enumerate(doc):
                    t = page.get_text("text").strip()
                    if t:
                        pages_text.append(f"--- Page {i+1} ---\n{t}")
                    else:
                        # Attempt to render page to image for EasyOCR if text layer is empty
                        try:
                            pix = page.get_pixmap(dpi=150)
                            img_bytes = pix.tobytes("png")
                            ocr_img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                            from src.image_analytics.vision_assistant import get_vision_assistant
                            assistant = get_vision_assistant()
                            ocr_txt = assistant._extract_text_ocr(ocr_img).strip()
                            if ocr_txt:
                                pages_text.append(f"--- Page {i+1} (OCR Extracted) ---\n{ocr_txt}")
                        except Exception:
                            pass
                full_pdf_text = "\n\n".join(pages_text).strip()
            except Exception as e:
                full_pdf_text = ""

        # Priority B: pypdf fallback
        if not full_pdf_text and pypdf:
            try:
                reader = pypdf.PdfReader(io.BytesIO(file_bytes))
                text_pages = []
                for i, page in enumerate(reader.pages):
                    p_text = page.extract_text() or ""
                    if p_text.strip():
                        text_pages.append(f"--- Page {i+1} ---\n{p_text.strip()}")
                full_pdf_text = "\n\n".join(text_pages).strip()
            except Exception:
                pass

        if full_pdf_text:
            return None, f"[PDF Document Content ({filename})]:\n{full_pdf_text[:4500]}"
        else:
            return None, f"[PDF Document ({filename})]: Uploaded PDF file contains no readable text pages or contains un-extracted scanned images."

    # 3. Word Documents (.docx, .doc)
    if ext in {".docx", ".doc"}:
        if docx:
            try:
                doc = docx.Document(io.BytesIO(file_bytes))
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                tables_text = []
                for t in doc.tables:
                    for row in t.rows:
                        tables_text.append(" | ".join(cell.text.strip() for cell in row.cells))
                full_doc = "\n".join(paragraphs)
                if tables_text:
                    full_doc += "\n\n[Tables]:\n" + "\n".join(tables_text)
                return None, f"[Word Document Content ({filename})]:\n{full_doc[:4000]}"
            except Exception as e:
                return None, f"[Word Document Error]: Could not read docx ({e})."

    # 4. PowerPoint Presentations (.pptx, .ppt)
    if ext in {".pptx", ".ppt"}:
        if pptx:
            try:
                prs = pptx.Presentation(io.BytesIO(file_bytes))
                slide_texts = []
                for i, slide in enumerate(prs.slides):
                    st = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            st.append(shape.text.strip())
                    if st:
                        slide_texts.append(f"--- Slide {i+1} ---\n" + "\n".join(st))
                full_ppt = "\n\n".join(slide_texts)
                return None, f"[PowerPoint Presentation Content ({filename}) - {len(prs.slides)} Slides]:\n{full_ppt[:4000]}"
            except Exception as e:
                return None, f"[PowerPoint Error]: Could not read pptx ({e})."

    # 5. Excel Spreadsheets & CSV (.xlsx, .xls, .csv)
    if ext in {".xlsx", ".xls", ".csv"}:
        if ext == ".csv":
            try:
                csv_text = file_bytes.decode("utf-8", errors="ignore")
                return None, f"[CSV Data Table ({filename})]:\n{csv_text[:3000]}"
            except Exception:
                pass
        if pd:
            try:
                df = pd.read_excel(io.BytesIO(file_bytes))
                summary = f"Columns: {list(df.columns)}\nTotal Rows: {len(df)}\n\nSample Data:\n{df.head(15).to_string()}"
                return None, f"[Excel Spreadsheet Data ({filename})]:\n{summary[:3500]}"
            except Exception as e:
                return None, f"[Excel Read Warning]: {e}"

    # 6. Video Files (.mp4, .avi, .mov, .mkv, .webm)
    if ext in {".mp4", ".avi", ".mov", ".mkv", ".webm"} or "video/" in content_type:
        if cv2:
            try:
                temp_video_path = f"temp_upload_{filename}"
                with open(temp_video_path, "wb") as f:
                    f.write(file_bytes)

                cap = cv2.VideoCapture(temp_video_path)
                total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
                fps = cap.get(cv2.CAP_PROP_FPS) or 30.0
                duration_sec = total_frames / fps if total_frames > 0 else 0

                # Grab frame at 50% midpoint of video
                cap.set(cv2.CAP_PROP_POS_FRAMES, int(total_frames * 0.5))
                ret, frame = cap.read()
                cap.release()

                if os.path.exists(temp_video_path):
                    os.remove(temp_video_path)

                if ret and frame is not None:
                    rgb_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                    pil_img = Image.fromarray(rgb_frame)
                    video_meta = f"[Video File Analysis ({filename}) - Duration: {duration_sec:.1f}s, Total Frames: {total_frames}, Frame Sampled at Midpoint]"
                    return pil_img, video_meta
            except Exception as e:
                return None, f"[Video Analysis Notice]: Duration/Frame extraction ({e})."

    # 7. Text & Code Files ONLY (Strict Extension Check)
    if ext in TEXT_EXTENSIONS or "text/" in content_type:
        try:
            raw_text = file_bytes.decode("utf-8", errors="ignore").strip()
            if raw_text:
                return None, f"[Text/Code File Content ({filename})]:\n{raw_text[:4000]}"
        except Exception:
            pass

    return None, f"[Uploaded File Metadata]: Filename '{filename}', size {len(file_bytes)} bytes."

